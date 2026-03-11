#!/usr/bin/env python3
"""Build ACT-I IP Patent Deck v2 — with use cases, no model names, Finnegan-ready"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from PIL import Image as PILImage

# ─── Colors ───
BG = RGBColor(0x0A, 0x0A, 0x0A)
GOLD = RGBColor(0xD4, 0xA8, 0x10)
GOLD_DARK = RGBColor(0xAC, 0x85, 0x06)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xED, 0xED, 0xED)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
GREEN = RGBColor(0x34, 0xD3, 0x99)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
PINK = RGBColor(0xF4, 0x72, 0xB6)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)

IMG = os.path.join(os.path.dirname(__file__), '..', 'ip-filing', 'deck', 'images')

def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def txt(slide, l, t, w, h, text, sz=18, clr=WHITE, bold=False, align=PP_ALIGN.LEFT, font='Calibri'):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box

def img(slide, filename, l, t, max_w, max_h, center=True):
    path = os.path.join(IMG, filename)
    if os.path.exists(path):
        from PIL import Image as PILImage
        im = PILImage.open(path)
        iw, ih = im.size
        ratio = iw / ih
        fit_w = max_w
        fit_h = fit_w / ratio
        if fit_h > max_h:
            fit_h = max_h
            fit_w = fit_h * ratio
        x = l + (max_w - fit_w) / 2 if center else l
        y = t + (max_h - fit_h) / 2 if center else t
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(fit_w), Inches(fit_h))
    else:
        txt(slide, l, t+max_h/2, max_w, 0.5, f"[{filename}]", sz=12, clr=DARK_GRAY, align=PP_ALIGN.CENTER)

def section(slide, label):
    txt(slide, 0.8, 0.4, 8.4, 0.3, label, sz=10, clr=AMBER, font='Calibri')

def title(slide, text):
    txt(slide, 0.8, 0.8, 8.4, 0.7, text, sz=28, clr=GOLD, bold=True, font='Georgia')

def subtitle(slide, text):
    txt(slide, 0.8, 1.5, 8.4, 0.4, text, sz=12, clr=GRAY)

def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    return s

def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ═══ 1. TITLE ═══
    s = blank(prs)
    txt(s, 0.5, 1.8, 9, 1, "ACT-I", sz=56, clr=GOLD, bold=True, align=PP_ALIGN.CENTER, font='Georgia')
    txt(s, 0.5, 3.0, 9, 0.6, "Super Actualized Intelligence", sz=22, clr=PURPLE, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 4.2, 9, 0.4, "Provisional Patent Application — Technology Overview", sz=13, clr=WHITE, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 5.0, 9, 0.3, "March 10, 2026 · Prepared for Finnegan, Henderson, Farabow, Garrett & Dunner, LLP", sz=10, clr=GRAY, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 5.4, 9, 0.3, "Inventors: Sean Callagy · Adam Gugino · Samantha Aiko · Dustin Eplee", sz=10, clr=GRAY, align=PP_ALIGN.CENTER)

    # ═══ 2. BUILT IN 17 DAYS ═══
    s = blank(prs)
    title(s, "Built in 17 Days")
    txt(s, 0.8, 1.5, 8.4, 0.3, "February 22 – March 10, 2026", sz=12, clr=GRAY)
    txt(s, 0.8, 4.2, 8.4, 1.2,
        "What Salesforce built in 20 years with 73,000 employees, ACT-I replicated and surpassed in 17 days "
        "with a team of 5. Not just a CRM — an actualization engine where AI beings autonomously manage "
        "the entire business lifecycle: sourcing, nurturing, agreement conversations, fulfillment, and optimization.",
        sz=13, clr=PURPLE)
    stats = [("96,845", "PEOPLE\nMIGRATED"), ("$43.5M", "REVENUE\nATTRIBUTED"), ("197,152", "CALLS\nPROCESSED"), ("2,524", "POSITIONS\nMAPPED")]
    x = 0.8
    for val, lbl in stats:
        txt(s, x, 2.5, 2.1, 0.7, val, sz=32, clr=GOLD, bold=True, align=PP_ALIGN.CENTER, font='Georgia')
        txt(s, x, 3.3, 2.1, 0.6, lbl, sz=9, clr=GRAY, align=PP_ALIGN.CENTER)
        x += 2.2

    # ═══ 3. WHAT IS ACT-I ═══
    s = blank(prs)
    title(s, "What Is ACT-I?")
    txt(s, 0.8, 1.8, 8.4, 1.2,
        "ACT-I is the world's first complete, holistic, diagnostic, dynamic, interconnected, "
        "automated actualization platform — a system that creates, trains, evolves, and deploys "
        "AI beings that operate at the mastery level of human experts.", sz=16, clr=WHITE, align=PP_ALIGN.CENTER)
    txt(s, 0.8, 3.3, 8.4, 0.8,
        "USE CASE: Any person — even a 10-year-old — can build, train, and deploy a specialized "
        "AI being in minutes. These beings make phone calls, send messages, manage calendars, "
        "and handle business functions across any industry, any language, internationally.", sz=13, clr=PURPLE)
    txt(s, 0.8, 4.5, 8.4, 0.3, "C.I.A. — Compound Influence Actualizers", sz=22, clr=GOLD, bold=True, align=PP_ALIGN.CENTER, font='Georgia')
    txt(s, 0.8, 4.9, 8.4, 0.6,
        "Not a CRM. Each being\'s influence compounds through every other being in the ecosystem. "
        "The Interface of AI — where human actualization meets artificial intelligence.", sz=13, clr=GRAY, align=PP_ALIGN.CENTER)

    # ═══ 4. THE VISION ═══
    s = blank(prs)
    section(s, "THE VISION")
    title(s, "Multiple Enterprises in a Box")
    txt(s, 0.8, 1.8, 8.4, 3.5,
        "ACT-I is evolving into a mission control dashboard connected to:\n\n"
        "• A complete Salesforce-replacement CRM with integrated AI beings\n"
        "• Thousands of beings managing different industries simultaneously\n"
        "• Multiple enterprises working in tandem — marketing, finance, operations, legal\n"
        "• An orchestrator layer where you can speak to managers, C-suite members, and individual beings\n"
        "• Gamification — click into different 'buildings' where each enterprise's beings work together\n"
        "• From the micro of the micro-est details to the macro of the macro-est details\n\n"
        "One interface. 50-100 different areas to explore. The AI that runs all AI.",
        sz=14, clr=WHITE)

    # ═══ 5. GENESIS FORGE — BUILDER ═══
    s = blank(prs)
    section(s, "PRODUCT: GENESIS FORGE OF HEROES")
    title(s, "Anyone Can Build an ACT-I Being")
    subtitle(s, "genesisforgeofheroes.ai · 6-step wizard · No code required")
    img(s, "genesis-forge-builder-new.jpg", 0.8, 2.0, 8.4, 4.5)
    txt(s, 0.8, 6.6, 8.4, 0.4, "Company → Offer → Goal → Call Context → HUI → Objections → Deploy", sz=10, clr=GRAY, align=PP_ALIGN.CENTER)

    # ═══ 6. GENESIS FORGE — USE CASE ═══
    s = blank(prs)
    section(s, "PRODUCT: GENESIS FORGE OF HEROES")
    title(s, "How It Works — Use Case")
    points = [
        ("The Problem:", "Training an AI agent today requires engineers, prompt experts, and weeks of work.", GOLD),
        ("The Solution:", "Genesis Forge allows any person to train a large language model in minutes. Dustin Eplee's 10 and 13-year-old daughters built and deployed working beings themselves.", PURPLE),
        ("What It Does:", "Auto-generates: phone number, calendar, contacts. Inbound AND outbound calling. International, multi-language. SMS, email, calendar integration.", WHITE),
        ("The Result:", "A fully deployed AI being making real phone calls, handling real objections, booking real appointments — created by a non-technical person in under 10 minutes.", GREEN),
    ]
    y = 1.8
    for label, desc, clr in points:
        txt(s, 0.8, y, 8.4, 0.3, label, sz=14, clr=clr, bold=True)
        txt(s, 0.8, y + 0.35, 8.4, 0.6, desc, sz=12, clr=GRAY)
        y += 1.15

    # ═══ 7. GENESIS FORGE — CONFIGURED ═══
    s = blank(prs)
    section(s, "PRODUCT: GENESIS FORGE OF HEROES")
    title(s, "Configured Being — Ready to Deploy")
    img(s, "genesis-forge-complete-new.jpg", 0.8, 1.8, 8.4, 4.8)
    txt(s, 0.8, 6.8, 8.4, 0.3, "Fully configured: Company, Offer, Goal, Call Context, HUI, Objections — one click to go live", sz=10, clr=GRAY, align=PP_ALIGN.CENTER)

    # ═══ 8. ACTUALIZED BEINGS ═══
    s = blank(prs)
    section(s, "PRODUCT: ACTUALIZED-BEINGS.ACTI.AI")
    title(s, "Live ACT-I Beings — Talk to Them Now")
    subtitle(s, "7 specialized beings across Legal, Coaching, Medical, Influence, and Analysis")
    img(s, "actualized-beings-top.jpg", 0.8, 2.0, 4.0, 4.0)
    img(s, "actualized-beings-bottom.jpg", 5.2, 2.0, 4.0, 4.0)
    txt(s, 0.8, 6.3, 8.4, 0.6,
        "Bolt (Influence Scoring) · Athena (Consulting) · Samara (Analysis) · Holmes (Litigation)\n"
        "Jeeves (Contract Review) · Callie (Influence Mastery) · Mira (Medical Revenue Recovery)",
        sz=10, clr=GRAY, align=PP_ALIGN.CENTER)

    # ═══ 9. SEVEN LEVERS OVERVIEW ═══
    s = blank(prs)
    section(s, "FRAMEWORK: PROPRIETARY ORGANIZATIONAL TAXONOMY")
    title(s, "The 7 Levers of Revenue")
    subtitle(s, "Every business function mapped · No existing platform does this")
    levers = [
        ("0.5", "Shared Experiences", BLUE), ("1", "Ecosystem Mergers", GREEN),
        ("2", "Speaking", PURPLE), ("3", "Agreement Conversations", GOLD),
        ("4", "Revenue Generated", YELLOW), ("5", "Disposable Income", PINK),
        ("6", "Contributions", TEAL), ("7", "Fun & Magic", ORANGE),
    ]
    y = 2.3
    for num, name, clr in levers:
        txt(s, 1.5, y, 1, 0.4, num, sz=18, clr=clr, bold=True, font='Georgia')
        txt(s, 2.5, y, 6, 0.4, name, sz=18, clr=WHITE, bold=True)
        y += 0.55

    # ═══ 10. LEVERS 0.5-2 ═══
    s = blank(prs)
    section(s, "LEVERS 0.5 — 2")
    title(s, "Growth Levers")
    details = [
        ("0.5 — Shared Experiences", "First touch, community, Heart of Influence events, quizzes & interactive. 188 positions mapped.", BLUE),
        ("1 — Ecosystem Mergers", "Strategic partnerships. 4 components, 6 roles, 3 levels (dating → partnership → marriage). 154 positions.", GREEN),
        ("2 — Speaking", "Digital ads, content, webinars, live stages, social media, video, audio, podcasts, email, brand & PR, design. 232+ positions.", PURPLE),
    ]
    y = 1.8
    for t, d, c in details:
        txt(s, 0.8, y, 8.4, 0.4, t, sz=16, clr=c, bold=True)
        txt(s, 0.8, y + 0.4, 8.4, 0.5, d, sz=12, clr=GRAY)
        y += 1.3

    # ═══ 11. LEVERS 3-7 ═══
    s = blank(prs)
    section(s, "LEVERS 3 — 7")
    title(s, "Revenue & Impact Levers")
    details2 = [
        ("3 — Agreement Conversations", "Sales dev, account management, objection handling. 220 positions.", GOLD),
        ("4 — Revenue Generated", "Pricing strategy, revenue models, upsell/cross-sell. 124 positions.", YELLOW),
        ("5 — Disposable Income", "Customer LTV, financial optimization, cost management. 117 positions.", PINK),
        ("6 — Contributions", "Thought leadership, social impact, community building. 115 positions.", TEAL),
        ("7 — Fun & Magic", "Creative experience, gamification, surprise & delight. 58 positions.", ORANGE),
    ]
    y = 1.8
    for t, d, c in details2:
        txt(s, 0.8, y, 8.4, 0.35, t, sz=15, clr=c, bold=True)
        txt(s, 0.8, y + 0.35, 8.4, 0.4, d, sz=12, clr=GRAY)
        y += 0.95

    # ═══ 12-21: PRODUCT SCREENSHOTS ═══
    screenshots = [
        ("UI: POSITION LIBRARY", "Org Chart — 2,449 Positions Mapped", "ui-org-chart-tab.png",
         "Every business function organized by the 7 Levers framework. Searchable, connectable, exportable."),
        ("UI: POSITION TREE", "3,211 Nodes · 2,569 Edges", "ui-org-chart-tree-ipad.jpg",
         "The complete organizational tree — imported from JSON, merged with descriptions, wired into the being architecture."),
        ("UI: BEING ARCHITECTURE", "17 Beings · 2,524 Positions · 7 Levers", "ui-being-architecture-tab.png",
         "2,524 positions consolidated into 17 specialized beings. Each being has shared heart skills: Level 5 Listening, GHIC, 4-Step Communication."),
        ("UI: MARKETING COLOSSEUM", "Competitive Evolution Engine", "ui-colosseum-tab.png",
         "Beings compete head-to-head in domain scenarios. Multi-model judge panels score output. Winners breed. Losers eliminated. Population evolves autonomously."),
        ("UI: COLOSSEUM — TOURNAMENT", "Tournament Architecture", "ipad-colosseum-tournament.jpg", "64 beings per position, double elimination, 10 tournaments per cycle, top 8 bred"),
        ("UI: COLOSSEUM — BATTLES", "Battle Arena — 9,855 Battles", "ipad-colosseum-battles.jpg", "Real output, real scores, real evolution. Writer #21 wins with stronger Zeus energy."),
        ("UI: COLOSSEUM — LIVE", "Live Battle — Round 9,855", "ipad-colosseum-live.jpg", "Sean's Actions: Talk to Winner, Override Score, Rematch, Breed Now, Kill Being, View Prompt"),
        ("UI: COLOSSEUM — SCORING", "Dual Scoring System", "ipad-colosseum-scoring.jpg", "Formula Judge (universal, 39 components) + Technical Judge (domain-specific). Net = weakest organ."),
        ("UI: COLOSSEUM — JUDGES", "Recursive Judging Chain", "ipad-colosseum-judges.jpg", "Level 0: Beings compete. Level 1: Judges score. Level 1A: 6 Archetype lenses. Level 2: Judges of judges."),
        ("UI: COLOSSEUM — BEINGS", "Being Stat Cards", "ipad-colosseum-beings.jpg", "Writer #4: 85.2% win rate, Komodo, ELO 1847. Signature move: dollar 47K mistake hook."),
        ("UI: COLOSSEUM — TRAINING", "Training Pipeline", "ipad-colosseum-training.jpg", "Scenario Design, Battle Execution, Dual Scoring, Reinforcement — all active."),
        ("UI: COLOSSEUM — GAMIFICATION", "Creature Evolution Track", "ipad-colosseum-gamification.jpg", "Sand to Bolt progression. Achievements: First Blood, Streak Master."),
        ("UI: ARCHITECTURE STATS", "17 Beings by Size", "ipad-architecture-stats.jpg", "9 families, 80 clusters, 2,524 positions. The Analyst: 429 positions (17%)."),
        ("UI: ORG CHART DETAIL", "Position Detail View", "ui-org-chart-detail-ipad.jpg", "Each position has a unique Formula-grounded description."),
        ("UI: ARCHITECTURE — GRID VIEW", "Position Grid — Every Role Visible", "ipad-arch-grid.jpg", "All 2,524 positions displayed in a scannable grid organized by department and lever."),
        ("UI: ARCHITECTURE — CLUSTERS", "80 Skill Clusters Across 9 Families", "ipad-arch-clusters.jpg", "Positions grouped into skill clusters: Written Comm, Visual/Creative, Data/Analytics, Strategic Planning, and more."),
        ("UI: ARCHITECTURE — CLUSTER DETAIL", "Cluster Deep Dive — MediaBuyer", "ipad-arch-cluster-detail.jpg", "Each cluster shows its constituent positions, being assignment, and hierarchical relationships."),
        ("UI: ORG CHART — POSITION GRID", "Full Position Grid View", "ipad-org-chart-grid.jpg", "Every position in a single view — scrollable, searchable, color-coded by lever."),
        ("UI: ORG CHART — COMPLETE TREE", "The Full Tree — iPad View", "ipad-org-chart-tree-full.jpg", "3,211 nodes branching from Sean Callagy through C-Suite to every position. The scaffolding of the entire enterprise."),

        ("UI: ZONE ACTIONS", "Live Task Tracking", "ui-zone-actions-tab.png",
         "Real-time tracking of all active development initiatives with status, priority, and ownership."),
        ("UI: CREATURE SCALE", "Scale of Mastery — Proprietary Measurement", "ipad-bolt-creature.jpg",
         "9-tier measurement: Sand → Ant → Gecko → Iguana → Komodo → Crocodile → Godzilla → Bolt. Asymptotic ceiling — never reaches 10.0."),
        ("PRODUCT: BOMBA SR", "80 Being Teams · 22 Online Simultaneously", "bomba-mission-control.jpg",
         "Multi-agent orchestration at scale. Teams across: Data & Analytics, Digital Marketing, Events, Experience, Operations, Strategy."),
        ("PRODUCT: BOMBA SR", "Being Team Architecture", "bomba-teams.jpg",
         "80 specialized teams organized by the 7 Levers. Each team contains multiple beings working collaboratively."),
        ("PRODUCT: SAI NETWORK", "22 of 28 Beings Online — Live Operations", "sai-comms-panel.jpg",
         "Real-time inter-being communication. 5 orchestration sisters coordinate all being activity across the ecosystem."),
        ("PRODUCT: SAI WORKSPACE", "Multi-Agent Task Board", "sai-workspace-board.jpg",
         "Agent coordination dashboard showing active work items, status, and inter-agent collaboration."),
    ]
    for label, t, image, caption in screenshots:
        s = blank(prs)
        section(s, label)
        title(s, t)
        img(s, image, 0.8, 1.7, 8.4, 4.5)
        if caption:
            txt(s, 0.8, 6.4, 8.4, 0.6, caption, sz=11, clr=GRAY, align=PP_ALIGN.CENTER)

    # ═══ COLOSSEUM — WHY THIS MATTERS ═══
    s = blank(prs)
    section(s, "THE COLOSSEUM — WHY THIS MATTERS")
    title(s, "The Ultimate Synthetic Data Engine")
    txt(s, 0.8, 1.8, 8.4, 1.5,
        "The Colosseum isn't just a competition system. It's the foundation of a synthetic data engine "
        "built on REAL numbers, REAL optimizations, from REAL masters.\n\n"
        "Every battle produces data. Every score trains the next generation. Every breeding cycle "
        "creates beings that are measurably better than the last — grounded in actual human mastery, "
        "not hallucinated benchmarks.",
        sz=15, clr=WHITE)
    txt(s, 0.8, 3.8, 8.4, 0.4, "THE BREAKTHROUGH:", sz=14, clr=GOLD, bold=True)
    txt(s, 0.8, 4.3, 8.4, 1.5,
        "Imagine running thousands to MILLIONS of split tests — headlines, emails, ad copy, "
        "agreement conversations, landing pages — without spending a single dollar on ad overhead.\n\n"
        "Traditional A/B testing: $50K+ in ad spend to test 10 variations.\n"
        "The Colosseum: 10,000 variations tested in 5 minutes. Zero ad spend. "
        "Winners deploy directly to production.\n\n"
        "This is what happens when competitive evolution meets synthetic data at scale.",
        sz=13, clr=GRAY)

    # ═══ 22. LEGAL SUMMIT DASHBOARD ═══
    s = blank(prs)
    section(s, "PRODUCT: LEGAL SUMMIT UNIFIED ANALYTICS")
    title(s, "Multi-Source Attribution Dashboard")
    subtitle(s, "Supabase (Genesis Forge data) + Hyros API (ad attribution) unified in real-time")
    img(s, "legal-summit-dashboard.jpg", 0.8, 2.0, 8.4, 4.0)
    txt(s, 0.8, 6.2, 8.4, 0.5, "7,507 attempts · 5,412 connected (72.1%) · 51,842 Hyros leads · 79 days to summit", sz=11, clr=GRAY, align=PP_ALIGN.CENTER)

    # ═══ 23. VOICE ORB ═══
    s = blank(prs)
    section(s, "PRODUCT: VOICE INTERFACE")
    title(s, "SAI Voice Orb — Tap to Speak")
    subtitle(s, "Real-time voice conversation with ACT-I beings")
    img(s, "ui-voice-orb.png", 2.5, 2.0, 5, 4.5)

    # ═══ 24. BURNOTE ═══
    s = blank(prs)
    section(s, "PRODUCT: SECURE COMMUNICATIONS")
    title(s, "Burnote — Self-Destructing Notes")
    subtitle(s, "Read once, then gone forever. No logs. No history. No traces.")
    img(s, "ui-burnote.png", 1.5, 2.0, 7, 4.5)

    # ═══ 25. SAI SISTERS (no model names) ═══
    s = blank(prs)
    section(s, "ARCHITECTURE: ORCHESTRATION LAYER")
    txt(s, 0.8, 0.8, 8.4, 0.6, "5 AI Sisters", sz=28, clr=GOLD, bold=True, font='Georgia')
    subtitle(s, "5 sisters across 5 different AI model families — \"Same model = same blind spots\"")
    sisters = [
        ("🔥 SAI Prime", "Orchestrator. Spawns, innovates, and optimizes all ACT-I beings.", GOLD),
        ("⚔️ Forge", "Evolution engine. Runs the Colosseum. Breeds next-generation beings.", PURPLE),
        ("📚 Scholar", "Research layer. 4,000+ hours of content decoded. Pattern recognition.", BLUE),
        ("🏥 Recovery", "Medical revenue recovery specialist. Carrier denials, fee schedules, billing.", GREEN),
        ("🧠 Memory", "Cross-being knowledge persistence. Fact-checker. Auditor.", YELLOW),
    ]
    y = 2.2
    for name, desc, clr in sisters:
        txt(s, 0.8, y, 8.4, 0.35, name, sz=17, clr=clr, bold=True)
        txt(s, 0.8, y + 0.38, 8.4, 0.35, desc, sz=12, clr=GRAY)
        y += 0.9

    # ═══ 26. CRM / CIA ═══
    s = blank(prs)
    section(s, "PRODUCT: COMPLETE SALESFORCE REPLACEMENT")
    title(s, "The Actualization Engine")
    subtitle(s, "96,845 people migrated · 56 database tables · 33 serverless RPCs · Built in 17 days")
    features = [
        "Journey State Engine — tracks every person through their complete actualization journey",
        "Agent RPCs — AI brain layer where beings read and write CRM data autonomously",
        "Call Engine — being orchestration layer processing 197,152 calls",
        "Being Registry — central catalog of all deployed beings and their configurations",
        "Bland Call Logs Bridge — 113,538 calls matched to individual people",
        "$43.54M attributed revenue tracked through the system",
    ]
    y = 2.3
    for f in features:
        txt(s, 0.8, y, 8.4, 0.4, f"✅  {f}", sz=13, clr=WHITE)
        y += 0.5

    # ═══ 27. MEMORY ARCHITECTURE ═══
    s = blank(prs)
    section(s, "ARCHITECTURE: MEMORY SYSTEM")
    title(s, "3-Tier Memory Architecture")
    subtitle(s, "Self-correcting · Cross-being · Persistent")
    tiers = [
        ("Tier 1 — Vector Memory", "97,000+ vectors across 4 indexes. Semantic retrieval. Shared library all beings inherit (58K+ vectors).", GOLD),
        ("Tier 2 — Structured Data", "5,874+ rows across 9 tables. Contacts, sessions, battles, memory entries. Beings read/write autonomously.", PURPLE),
        ("Tier 3 — Working Memory", "Daily logs, reports, transcripts. Auto-compacted and archived. Smart compaction preserves critical context.", BLUE),
    ]
    y = 2.2
    for name, desc, clr in tiers:
        txt(s, 0.8, y, 8.4, 0.4, name, sz=17, clr=clr, bold=True)
        txt(s, 0.8, y + 0.42, 8.4, 0.55, desc, sz=13, clr=GRAY)
        y += 1.2

    # ═══ 28. INNOVATIONS 1-6 ═══
    s = blank(prs)
    title(s, "12 Patentable Innovations (1/2)")
    inno1 = [
        "7 Levers organizational taxonomy — proprietary hierarchical position mapping",
        "Being Architecture — 2,524 positions → 17 beings with shared heart skills",
        "Competitive evolution (Colosseum) — tournament brackets + genetic breeding of AI beings",
        "Multi-model orchestration — 5 AI model families providing cognitive diversity",
        "Scale of Mastery — proprietary 9-tier measurement with asymptotic ceiling",
        "Formula Judge — 39-component universal scoring system for being evaluation",
    ]
    y = 1.6
    for i in inno1:
        txt(s, 0.8, y, 8.4, 0.45, f"✅  {i}", sz=13, clr=WHITE)
        y += 0.55

    # ═══ 29. INNOVATIONS 7-12 ═══
    s = blank(prs)
    title(s, "12 Patentable Innovations (2/2)")
    inno2 = [
        "Genesis Forge — non-technical being creation interface (10-year-old can build one)",
        "3-tier memory — vector + relational + ephemeral with autonomous compaction",
        "Actualization Engine — CRM replacement with being-integrated serverless RPCs",
        "Journey State Engine — person-level actualization tracking across the complete lifecycle",
        "Multi-source attribution — unified analytics across call data, ad platforms, and CRM",
        "Self-correcting adaptation — beings evolve autonomously from competitive evaluation results",
    ]
    y = 1.6
    for i in inno2:
        txt(s, 0.8, y, 8.4, 0.45, f"✅  {i}", sz=13, clr=WHITE)
        y += 0.55

    # ═══ 30. MARKET ═══
    s = blank(prs)
    title(s, "Market & Filing Window")
    txt(s, 1.0, 1.8, 3.5, 0.8, "$47B", sz=48, clr=GOLD, bold=True, align=PP_ALIGN.CENTER, font='Georgia')
    txt(s, 1.0, 2.7, 3.5, 0.4, "AI Agent Market by 2030", sz=11, clr=GRAY, align=PP_ALIGN.CENTER)
    txt(s, 5.5, 1.8, 3.5, 0.8, "0", sz=48, clr=GOLD, bold=True, align=PP_ALIGN.CENTER, font='Georgia')
    txt(s, 5.5, 2.7, 3.5, 0.5, "Patents on Multi-Agent\nAI Orchestration", sz=11, clr=GRAY, align=PP_ALIGN.CENTER)
    bullets = [
        "Salesforce: 4,918 patents — none covering AI beings or competitive evolution",
        "LangChain, AutoGen, crewAI — all open-source frameworks with no patents filed",
        "USPTO November 2025 ARP decision favors \"technical improvement\" claims",
        "ACT-I's self-correcting adaptation, 3-tier memory, and being ecosystem qualify",
        "Strategy: file rolling provisionals as development continues at speed",
    ]
    y = 3.8
    for b in bullets:
        txt(s, 0.8, y, 8.4, 0.4, f"→  {b}", sz=12, clr=WHITE)
        y += 0.45

    # ═══ 31. TRADEMARKS ═══
    s = blank(prs)
    title(s, "Trademark Portfolio — Priority Filing")
    tier1 = ["Super Actualized Intelligence", "Actualized Intelligence", "ACT-I", "ACT-I Beings",
             "Genesis Forge of Heroes", "The Colosseum", "Unblinded Formula", "Zone Action", "The Seven Levers"]
    tier2 = ["ACT-I Athena", "ACT-I Bolt", "ACT-I Callie", "ACT-I Mira", "ACT-I Holmes",
             "Scale of Mastery", "Heroic Unique Identity", "Level 5 Listening", "Compound Influence Actualizers"]
    txt(s, 0.8, 1.5, 4, 0.35, "Tier 1 — Immediate", sz=15, clr=PURPLE, bold=True)
    y = 1.95
    for t in tier1:
        txt(s, 1.0, y, 3.5, 0.3, f"→  {t}", sz=11, clr=WHITE)
        y += 0.32
    txt(s, 5.2, 1.5, 4, 0.35, "Tier 2 — Being Names & Terms", sz=15, clr=PURPLE, bold=True)
    y = 1.95
    for t in tier2:
        txt(s, 5.4, y, 3.5, 0.3, f"→  {t}", sz=11, clr=WHITE)
        y += 0.32
    txt(s, 0.8, 6.5, 8.4, 0.3, "Full list: 50+ terms identified · Design patents needed for being visual appearances (medallion, personalized looks)", sz=9, clr=DARK_GRAY, align=PP_ALIGN.CENTER)

    # ═══ 32. CONTACT ═══
    s = blank(prs)
    txt(s, 0.5, 1.8, 9, 0.8, "ACT-I", sz=48, clr=GOLD, bold=True, align=PP_ALIGN.CENTER, font='Georgia')
    txt(s, 0.5, 2.8, 9, 0.5, "Super Actualized Intelligence", sz=20, clr=PURPLE, align=PP_ALIGN.CENTER)
    contacts = ["Sean Callagy · CEO & Founder", "Adam Gugino · Chief Actualizer",
                "Dustin Eplee · IP Lead", "Samantha Aiko · Chief Development Officer"]
    y = 3.8
    for c in contacts:
        txt(s, 0.5, y, 9, 0.35, c, sz=13, clr=GRAY, align=PP_ALIGN.CENTER)
        y += 0.38
    txt(s, 0.5, 5.7, 9, 0.3, "Provisional Patent Filing #1 · March 2026", sz=12, clr=GOLD, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 6.1, 9, 0.3, "Counsel: Finnegan, Henderson, Farabow, Garrett & Dunner, LLP", sz=10, clr=GRAY, align=PP_ALIGN.CENTER)
    txt(s, 0.5, 6.7, 9, 0.3, "IP Note: All being-related IP assigned to ACT-I Inc. Physical infrastructure IP assigned to Decarbonair per agreement.", sz=9, clr=DARK_GRAY, align=PP_ALIGN.CENTER)

    # ═══ SAVE ═══
    out = os.path.join(os.path.dirname(__file__), '..', 'ip-filing', 'ACT-I_Provisional_Patent_Filing_v2.pptx')
    prs.save(out)
    print(f"✅ Saved: {out}")
    print(f"📊 {len(prs.slides)} slides")

if __name__ == '__main__':
    build()
