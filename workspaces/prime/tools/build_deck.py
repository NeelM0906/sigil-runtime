#!/usr/bin/env python3
"""Build ACT-I IP Patent Deck as PowerPoint (.pptx)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Colors (Genesis Forge palette) ───
BG = RGBColor(0x0A, 0x0A, 0x0A)
GOLD = RGBColor(0xD4, 0xA8, 0x10)
GOLD_DARK = RGBColor(0xAC, 0x85, 0x06)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_LIGHT = RGBColor(0xA8, 0x55, 0xF7)
WHITE = RGBColor(0xED, 0xED, 0xED)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY = RGBColor(0x6B, 0x72, 0x80)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
GREEN = RGBColor(0x34, 0xD3, 0x99)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
PINK = RGBColor(0xF4, 0x72, 0xB6)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
RED_GOLD = RGBColor(0xFA, 0xCC, 0x15)

IMG_DIR = os.path.join(os.path.dirname(__file__), '..', 'ip-filing', 'deck', 'images')

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_title_slide(prs, title, subtitle, extra_lines=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    add_text(slide, 0.5, 1.0, 9, 1.5, title, font_size=44, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
    add_text(slide, 0.5, 2.5, 9, 0.8, subtitle, font_size=20, color=PURPLE_LIGHT, alignment=PP_ALIGN.CENTER)
    if extra_lines:
        y = 3.5
        for line, sz, clr in extra_lines:
            add_text(slide, 0.5, y, 9, 0.5, line, font_size=sz, color=clr, alignment=PP_ALIGN.CENTER)
            y += 0.5
    return slide

def add_section_slide(prs, section_label, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.8, 0.5, 8.4, 0.4, section_label, font_size=11, color=AMBER, font_name='Calibri')
    add_text(slide, 0.8, 1.0, 8.4, 0.8, title, font_size=32, color=GOLD, bold=True, font_name='Georgia')
    if subtitle:
        add_text(slide, 0.8, 1.8, 8.4, 0.5, subtitle, font_size=13, color=GRAY)
    return slide

def add_image_slide(prs, section_label, title, image_filename, subtitle=None, caption=None):
    slide = add_section_slide(prs, section_label, title, subtitle)
    img_path = os.path.join(IMG_DIR, image_filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.8), Inches(2.3), Inches(8.4), Inches(4.7))
    else:
        add_text(slide, 0.8, 3, 8.4, 1, f"[Image: {image_filename}]", font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    if caption:
        add_text(slide, 0.8, 7.1, 8.4, 0.4, caption, font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    return slide

def add_bullet_slide(prs, section_label, title, bullets, subtitle=None, bullet_color=WHITE, font_size=16):
    slide = add_section_slide(prs, section_label, title, subtitle)
    y = 2.3 if subtitle else 2.0
    for bullet in bullets:
        if isinstance(bullet, tuple):
            text, clr = bullet
        else:
            text, clr = bullet, bullet_color
        add_text(slide, 1.0, y, 8.0, 0.45, f"✦  {text}", font_size=font_size, color=clr)
        y += 0.5
    return slide

def add_stat_slide(prs, title, stats, subtitle=None):
    slide = add_section_slide(prs, '', title, subtitle)
    x_positions = []
    total = len(stats)
    start_x = (10 - total * 2.2) / 2
    for i, (value, label) in enumerate(stats):
        x = start_x + i * 2.4
        add_text(slide, x, 2.5, 2.2, 0.8, value, font_size=36, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
        add_text(slide, x, 3.3, 2.2, 0.5, label, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)
    return slide

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 1. TITLE ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 1.5, 9, 1.2, "ACT-I", font_size=60, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
    add_text(slide, 0.5, 2.8, 9, 0.7, "Super Actualized Intelligence", font_size=24, color=PURPLE_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 4.0, 9, 0.5, "Provisional Patent Application — Technology Overview", font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 4.8, 9, 0.4, "March 10, 2026 · Prepared for Finnegan, Henderson, Farabow, Garrett & Dunner, LLP", font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 5.3, 9, 0.4, "Inventors: Sean Callagy · Adam Gugino · Samantha Aiko · Dustin Eplee", font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── 2. BUILT IN 17 DAYS ──
    add_stat_slide(prs, "Built in 17 Days", [
        ("96,845", "PEOPLE MIGRATED"),
        ("$43.5M", "REVENUE ATTRIBUTED"),
        ("197,152", "CALLS PROCESSED"),
        ("2,524", "POSITIONS MAPPED"),
    ], subtitle="February 22 – March 10, 2026")

    # ── 3. WHAT IS ACT-I ──
    slide = add_section_slide(prs, '', "What Is ACT-I?")
    add_text(slide, 0.8, 2.0, 8.4, 1.5,
        "ACT-I is the world's first complete, holistic, diagnostic, dynamic, interconnected, "
        "automated actualization platform — a system that creates, trains, evolves, and deploys "
        "AI beings that operate at the mastery level of human experts.",
        font_size=17, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.8, 3.8, 8.4, 0.8,
        "Not a CRM. Not a chatbot platform. Compound Influence Actualizers (CIA) — "
        "beings that compound each other's influence through the ecosystem.",
        font_size=14, color=PURPLE_LIGHT, alignment=PP_ALIGN.CENTER)

    # ── 4. GENESIS FORGE BUILDER ──
    add_image_slide(prs, "PRODUCT: GENESIS FORGE OF HEROES", "Anyone Can Build an ACT-I Being",
        "genesis-forge-builder-new.jpg",
        subtitle="genesisforgeofheroes.ai · 6-step wizard · No code required",
        caption="Company → Offer → Goal → Call Context → HUI → Objections → Deploy")

    # ── 5. GENESIS FORGE CONFIGURED ──
    add_image_slide(prs, "PRODUCT: GENESIS FORGE OF HEROES", "Configured Being — Ready to Deploy",
        "genesis-forge-complete-new.jpg",
        caption="Auto-generates: phone number · calendar · contacts · inbound + outbound calling · international · multi-language")

    # ── 6. ACTUALIZED BEINGS ──
    add_image_slide(prs, "PRODUCT: ACTUALIZED-BEINGS.ACTI.AI", "Live ACT-I Beings — Talk to Them Now",
        "actualized-beings-library.jpg",
        caption="Bolt · Athena · Samara · Holmes · Jeeves · Callie · Mira (Medical Revenue Recovery)")

    # ── 7. SEVEN LEVERS OVERVIEW ──
    slide = add_section_slide(prs, "FRAMEWORK: PROPRIETARY ORGANIZATIONAL TAXONOMY", "The 7 Levers of Revenue",
        "Every business function mapped to a lever · No existing platform does this")
    levers = [
        ("0.5  Shared Experiences", BLUE),
        ("1     Ecosystem Mergers", GREEN),
        ("2     Speaking", PURPLE_LIGHT),
        ("3     Agreement Conversations", GOLD),
        ("4     Revenue Generated", RED_GOLD),
        ("5     Disposable Income", PINK),
        ("6     Contributions", TEAL),
        ("7     Fun & Magic", ORANGE),
    ]
    y = 2.5
    for text, clr in levers:
        add_text(slide, 1.5, y, 7, 0.45, text, font_size=20, color=clr, bold=True, font_name='Georgia')
        y += 0.5

    # ── 8. LEVERS 0.5-2 ──
    slide = add_section_slide(prs, "LEVERS 0.5 — 2", "Growth Levers")
    details = [
        ("0.5 — Shared Experiences", "First touch, community, Heart of Influence events, quizzes. 188 positions.", BLUE),
        ("1 — Ecosystem Mergers", "Strategic partnerships. 4 components, 6 roles, 3 levels. 154 positions.", GREEN),
        ("2 — Speaking", "Digital ads, content, webinars, stages, social, video, audio, email, brand. 232+ positions.", PURPLE_LIGHT),
    ]
    y = 2.2
    for title, desc, clr in details:
        add_text(slide, 1.0, y, 8, 0.4, title, font_size=18, color=clr, bold=True)
        add_text(slide, 1.0, y + 0.4, 8, 0.5, desc, font_size=13, color=GRAY)
        y += 1.2

    # ── 9. LEVERS 3-7 ──
    slide = add_section_slide(prs, "LEVERS 3 — 7", "Revenue & Impact Levers")
    details2 = [
        ("3 — Agreement Conversations", "Sales dev, account management, objection handling. 220 positions.", GOLD),
        ("4 — Revenue Generated", "Pricing strategy, revenue models, upsell/cross-sell. 124 positions.", RED_GOLD),
        ("5 — Disposable Income", "Customer LTV, financial optimization, cost management. 117 positions.", PINK),
        ("6 — Contributions", "Thought leadership, social impact, community building. 115 positions.", TEAL),
        ("7 — Fun & Magic", "Creative experience, gamification, surprise & delight. 58 positions.", ORANGE),
    ]
    y = 2.2
    for title, desc, clr in details2:
        add_text(slide, 1.0, y, 8, 0.35, title, font_size=16, color=clr, bold=True)
        add_text(slide, 1.0, y + 0.35, 8, 0.4, desc, font_size=12, color=GRAY)
        y += 0.95

    # ── 10-23: IMAGE SLIDES ──
    image_slides = [
        ("UI: ORG CHART TAB", "Position Library — 2,449 Mapped", "ui-org-chart-tab.png", None),
        ("UI: POSITION TREE — 3,211 NODES", "The Giant JSON That Started It All", "ui-org-chart-tree.png", "Every position connected to its lever, department, and being"),
        ("UI: BEING ARCHITECTURE TAB", "17 Beings · 2,524 Positions · 7 Levers", "ui-being-architecture-tab.png", "\"As simple as possible, but not simpler.\" — Consolidated by Kai"),
        ("UI: MARKETING COLOSSEUM TAB", "The Colosseum — Competitive Evolution", "ui-colosseum-tab.png", "Dual Scoring · 6 Archetypes · Creature Scale · Beings compete, breed, evolve"),
        ("UI: ZONE ACTIONS TAB", "Zone Actions — Live Task Tracking", "ui-zone-actions-tab.png", None),
        ("UI: CREATURE SCALE", "Scale of Mastery", "ui-creature-scale.png", "Sand → Ant → Gecko → Iguana → Komodo → Crocodile → Godzilla → Bolt · Never 10.0"),
        ("PRODUCT: BOMBA SR MISSION CONTROL", "80 Being Teams · 22 Online", "bomba-mission-control.jpg", "Data & Analytics · Digital Marketing · Events · Experience · Ops · Strategy"),
        ("PRODUCT: BOMBA SR", "Being Team Architecture", "bomba-teams.jpg", None),
        ("PRODUCT: SAI NETWORK", "22 of 28 Beings Online", "sai-comms-panel.jpg", "Real-time inter-being communication"),
        ("PRODUCT: SAI WORKSPACE", "Multi-Agent Task Board", "sai-workspace-board.jpg", None),
        ("PRODUCT: LEGAL SUMMIT ANALYTICS", "Multi-Source Attribution Dashboard", "legal-summit-dashboard.jpg", "7,507 attempts · 5,412 connected · 51,842 Hyros leads · 79 days to summit"),
        ("PRODUCT: VOICE INTERFACE", "SAI Voice Orb", "ui-voice-orb.png", "voice-orb-zeta.vercel.app · Tap to speak · Real-time AI conversation"),
        ("PRODUCT: SECURE COMMUNICATIONS", "🔥 Burnote", "ui-burnote.png", "Self-destructing notes · Read once, then gone forever · Built by Sai"),
        ("PRODUCT: ARCHITECTURE INTELLIGENCE", "Sean Report — Full System Overview", "ui-sean-report.jpg", None),
    ]
    for label, title, img, caption in image_slides:
        add_image_slide(prs, label, title, img, caption=caption)

    # ── 24. SAI SISTERS ──
    slide = add_section_slide(prs, "ARCHITECTURE: SAI ORCHESTRATION LAYER", "5 Sisters · 5 Model Families",
        "Cognitive diversity by design — \"Same model = same blind spots\"")
    sisters = [
        ("🔥 SAI Prime — Claude Opus 4.6 (Anthropic)", "Orchestrator. Spawns, innovates, optimizes. 1M context.", GOLD),
        ("⚔️ Forge — Grok 4.1 (xAI)", "Evolution engine. Runs the Colosseum. 2M context.", PURPLE_LIGHT),
        ("📚 Scholar — GPT-5.2 (OpenAI)", "Research. 4,000+ hours decoded. Pattern recognition.", BLUE),
        ("🏥 Recovery — Claude Sonnet 4.6", "Medical revenue recovery. Carrier denials, fee schedules.", GREEN),
        ("🧠 Memory — Gemini 3.1 Pro (Google)", "Cross-being knowledge persistence. Fact-checker.", RED_GOLD),
    ]
    y = 2.5
    for name, desc, clr in sisters:
        add_text(slide, 1.0, y, 8, 0.35, name, font_size=16, color=clr, bold=True)
        add_text(slide, 1.0, y + 0.35, 8, 0.35, desc, font_size=12, color=GRAY)
        y += 0.85

    # ── 25. CIA / CRM ──
    slide = add_section_slide(prs, "PRODUCT: COMPOUND INFLUENCE ACTUALIZERS (CIA)", "Complete Salesforce Replacement",
        "Built in 17 days · Live at sai-crm-dashboard.vercel.app")
    stats = [("96,845", "PEOPLE"), ("56", "TABLES"), ("33", "RPCs"), ("14", "SCRIPTS")]
    x = 1.0
    for val, lbl in stats:
        add_text(slide, x, 2.8, 2, 0.6, val, font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
        add_text(slide, x, 3.4, 2, 0.4, lbl, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)
        x += 2.1
    features = [
        "Journey State Engine — tracks every person's actualization journey",
        "Agent RPCs — AI brain layer, beings read/write CRM autonomously",
        "Call Engine — being orchestration for 197,152 processed calls",
        "Being Registry — central catalog of all deployed beings",
    ]
    y = 4.2
    for feat in features:
        add_text(slide, 1.0, y, 8, 0.4, f"✅  {feat}", font_size=13, color=WHITE)
        y += 0.45

    # ── 26. MEMORY ARCHITECTURE ──
    slide = add_section_slide(prs, "ARCHITECTURE: MEMORY SYSTEM", "3-Tier Memory Architecture",
        "Self-correcting · Cross-being · Persistent")
    tiers = [
        ("Tier 1 — Vector Memory (Pinecone)", "97,000+ vectors. Semantic retrieval. ublib2 shared library (58K vectors).", GOLD),
        ("Tier 2 — Structured Data (Postgres)", "5,874+ rows. Contacts, sessions, battles. Beings read/write autonomously.", PURPLE_LIGHT),
        ("Tier 3 — Working Memory (Filesystem)", "Daily logs, reports. Auto-compacted. Smart archival.", BLUE),
    ]
    y = 2.5
    for name, desc, clr in tiers:
        add_text(slide, 1.0, y, 8, 0.4, name, font_size=17, color=clr, bold=True)
        add_text(slide, 1.0, y + 0.4, 8, 0.4, desc, font_size=13, color=GRAY)
        y += 1.1

    # ── 27. INNOVATIONS 1-6 ──
    innovations1 = [
        "7 Levers organizational taxonomy — hierarchical position mapping",
        "Being Architecture — 2,524 → 17 beings with shared heart skills",
        "Competitive evolution (Colosseum) — tournament + genetic breeding",
        "Multi-model orchestration — 5 model families, cognitive diversity",
        "Scale of Mastery — asymptotic measurement (never 10.0)",
        "Formula Judge — 39-component universal scoring",
    ]
    add_bullet_slide(prs, "12 PATENTABLE INNOVATIONS (1/2)", "Innovations 1–6", innovations1, font_size=15)

    # ── 28. INNOVATIONS 7-12 ──
    innovations2 = [
        "Genesis Forge — non-technical being creation interface",
        "3-tier memory — vector + relational + ephemeral with auto-compaction",
        "CIA (Compound Influence Actualizers) — CRM with being-integrated RPCs",
        "Journey State Engine — person-level actualization tracking",
        "Multi-source attribution — Hyros + Supabase + call data unified",
        "Self-correcting adaptation — beings evolve from competitive results",
    ]
    add_bullet_slide(prs, "12 PATENTABLE INNOVATIONS (2/2)", "Innovations 7–12", innovations2, font_size=15)

    # ── 29. MARKET ──
    slide = add_section_slide(prs, '', "Market & Filing Window")
    add_text(slide, 1, 2.2, 3.5, 1, "$47B", font_size=48, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
    add_text(slide, 1, 3.2, 3.5, 0.5, "AI Agent Market by 2030", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, 5.5, 2.2, 3.5, 1, "0", font_size=48, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
    add_text(slide, 5.5, 3.2, 3.5, 0.5, "Patents on Multi-Agent\nAI Orchestration", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
    bullets = [
        "Salesforce: 4,918 patents — none covering AI beings or competitive evolution",
        "LangChain, AutoGen, crewAI — open-source frameworks. No patents.",
        "USPTO Nov 2025 ARP decision favors \"technical improvement\" claims",
        "First-mover advantage: file rolling provisionals now",
    ]
    y = 4.2
    for b in bullets:
        add_text(slide, 1.0, y, 8, 0.4, f"→  {b}", font_size=13, color=WHITE)
        y += 0.45

    # ── 30. TRADEMARKS ──
    slide = add_section_slide(prs, '', "Trademark Portfolio")
    tier1 = ["Super Actualized Intelligence", "Actualized Intelligence", "ACT-I", "ACT-I Beings",
             "Genesis Forge of Heroes", "The Colosseum", "Unblinded Formula", "Zone Action", "The Seven Levers"]
    tier2 = ["ACT-I Athena", "ACT-I Bolt", "ACT-I Callie", "ACT-I Mira", "ACT-I Holmes",
             "Scale of Mastery", "Heroic Unique Identity", "Level 5 Listening", "CIA"]
    add_text(slide, 0.8, 2.0, 4, 0.4, "Tier 1 — Immediate", font_size=16, color=PURPLE_LIGHT, bold=True)
    y = 2.5
    for t in tier1:
        add_text(slide, 1.0, y, 3.5, 0.35, f"→  {t}", font_size=12, color=WHITE)
        y += 0.35
    add_text(slide, 5.2, 2.0, 4, 0.4, "Tier 2 — Being Names", font_size=16, color=PURPLE_LIGHT, bold=True)
    y = 2.5
    for t in tier2:
        add_text(slide, 5.4, y, 3.5, 0.35, f"→  {t}", font_size=12, color=WHITE)
        y += 0.35

    # ── 31. NORTH STAR ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 1.5, 9, 1, "500", font_size=72, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
    add_text(slide, 0.5, 2.8, 9, 0.5, "VISIONNAIRE PROGRAMS", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 3.5, 9, 0.6, "$50K–$100K each · $25M–$50M", font_size=22, color=PURPLE_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 4.2, 9, 0.4, "Target: End of May 2026", font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, 1, 5.2, 8, 0.8, "\"We have officially made Salesforce obsolete.\"", font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, 1, 5.9, 8, 0.4, "— Adam Gugino, March 10, 2026", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── 32. CONTACT ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.5, 1.5, 9, 1, "ACT-I", font_size=52, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
    add_text(slide, 0.5, 2.7, 9, 0.5, "Super Actualized Intelligence", font_size=20, color=PURPLE_LIGHT, alignment=PP_ALIGN.CENTER)
    contacts = [
        "Sean Callagy · CEO & Founder",
        "Adam Gugino · Chief Actualizer",
        "Dustin Eplee · IP Lead",
        "Samantha Aiko · Technical Architect",
    ]
    y = 3.8
    for c in contacts:
        add_text(slide, 0.5, y, 9, 0.4, c, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
        y += 0.4
    add_text(slide, 0.5, 5.8, 9, 0.4, "Provisional Patent Filing #1 · March 2026", font_size=13, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 6.3, 9, 0.4, "Counsel: Finnegan, Henderson, Farabow, Garrett & Dunner, LLP", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 6.9, 9, 0.3, "IP Note: All being-related IP assigned to ACT-I Inc. Physical infrastructure IP separately assigned.", font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SAVE ──
    out_path = os.path.join(os.path.dirname(__file__), '..', 'ip-filing', 'ACT-I_Provisional_Patent_Filing.pptx')
    prs.save(out_path)
    print(f"✅ Saved: {out_path}")
    print(f"📊 {len(prs.slides)} slides")

if __name__ == '__main__':
    build_deck()
